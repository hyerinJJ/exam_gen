from pathlib import Path


def read_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        pages.append(f"[페이지 {i}]\n{page.extract_text() or ''}")
    return "\n\n".join(pages)


def read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        slides.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _check_ffmpeg() -> None:
    import subprocess
    try:
        subprocess.run(["ffmpeg", "-version"], capture_output=True, check=False)
    except FileNotFoundError:
        raise RuntimeError(
            "ffmpeg가 설치되어 있지 않거나 PATH에 없습니다.\n"
            "설치 방법 (Windows): winget install ffmpeg\n"
            "설치 후 터미널을 재시작하고 다시 실행하세요."
        )


def read_video(path: Path) -> str:
    import shutil
    import tempfile
    import whisper

    _check_ffmpeg()
    model = whisper.load_model("base")

    # whisper/ffmpeg가 Windows에서 한글·공백 경로를 처리 못하므로
    # ASCII 이름의 임시 파일로 복사 후 transcribe
    suffix = path.suffix.lower()
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp_path = Path(tmp.name)

    try:
        shutil.copy2(str(path), str(tmp_path))
        result = model.transcribe(str(tmp_path))
    finally:
        tmp_path.unlink(missing_ok=True)

    return result["text"]


def read_file(file_path: str) -> str:
    path = Path(file_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {path}")

    ext = path.suffix.lower()
    if ext == ".pdf":
        return read_pdf(path)
    elif ext in (".pptx", ".ppt"):
        return read_pptx(path)
    elif ext in (".mp4", ".mov", ".avi", ".mkv"):
        return read_video(path)
    elif ext == ".txt":
        with open(str(path), encoding="utf-8") as f:
            return f.read()
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext} ({path})")
